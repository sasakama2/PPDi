from kivy.config import Config
Config.set('graphics', 'fullscreen', '0')
Config.set('graphics', 'width', '670')
Config.set('graphics', 'height', '480')
Config.write()
from kivy.core.text import LabelBase, DEFAULT_FONT  
from kivy.resources import resource_add_path
resource_add_path('c:/Windows/Fonts')  
LabelBase.register(DEFAULT_FONT, 'YuGothB.ttc')

from kivy.app import App
from kivy.core.text import DEFAULT_FONT, Label, LabelBase
from kivy.core.window import Window
from kivy.graphics import Color, Line
from kivy.lang import Builder
from kivy.properties import ObjectProperty, StringProperty, BooleanProperty
from kivy.resources import resource_add_path
from kivy.uix.boxlayout import BoxLayout
from kivy.uix.popup import Popup
from kivy.uix.screenmanager import (Screen, ScreenManager,
                                    ScreenManagerException)
from kivy.uix.behaviors import ButtonBehavior
from kivy.uix.image import Image as IMG
from kivy.uix.widget import Widget
from kivy.utils import get_color_from_hex
from pptx import Presentation
from pptx.dml.color import RGBColor,MSO_THEME_COLOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN, MSO_AUTO_SIZE
from pptx.util import Cm, Pt
from kivy.clock import Clock
from kivy.uix.checkbox import CheckBox

import pandas as pd
import win32com.client
import numpy as np
import os
import glob
from matplotlib import pyplot as plt
import shutil

import PIL.ImageDraw
from PIL import  Image, ImageDraw, ImageFont
from skimage import color,io
import scipy.stats
import re



class diagnosis(Screen):
    def __init__(self, **kwargs):
        super(diagnosis, self).__init__(**kwargs)
        pass

    def shape_color_pass(self,shp,new_zukei):
        '''
        修正後の図形の色を司る
        '''
    #枠線の色の認識・引き継ぎ
        try:
            line_color=shp.line.color.rgb
            new_zukei.line.color.rgb=RGBColor(int(line_color[-3]),int(line_color[-2]),int(line_color[-1]))
        except AttributeError:
            try:
                line_color=shp.line.color.theme_color
                line_color_bright=shp.line.color.brightness
                new_zukei.line.color.theme_color=line_color
                new_zukei.line.color.brightness=line_color_bright
            except AttributeError:
                new_zukei.line.color.rgb=RGBColor(65,113,156)
        #塗りつぶし色の認識・引き継ぎ＆枠線色の削除
        try:#標準の色
            shape_color=shp.fill.fore_color.rgb
            colo=RGBColor(int(shape_color[-3]),int(shape_color[-2]),int(shape_color[-1]))
            new_zukei.fill.solid()
            new_zukei.fill.fore_color.rgb=colo

            #ここに塗りつぶし色の変更（赤色やめるとか）
            
        except AttributeError:#テーマの色
            shape_color=shp.fill.fore_color.theme_color
            shape_color_bright=shp.fill.fore_color.brightness
            new_zukei.fill.solid()
            new_zukei.fill.fore_color.theme_color=shape_color
            new_zukei.fill.fore_color.brightness=shape_color_bright
    
        except TypeError:#元の図形がテーマスタイルor塗りつぶしが初期色orグラデーション
            new_zukei.fill.solid()
            new_zukei.fill.fore_color.theme_color=MSO_THEME_COLOR.ACCENT_1#初期設定の色=青アクセント１

    def text_color_pass(self,run,new_zukei):
        '''
        図形の文字色をひきつぐ関数。不適切な色は変更する
        '''
        try:#標準の色or設定されたRGBを認識
            text_color=run.font.color.rgb
            #print(run.text)
            #print(int(text_color[-3]),int(text_color[-2]),int(text_color[-1]))#RGB
            for pg2 in new_zukei.text_frame.paragraphs:
                for run2 in pg2.runs:
                    run2.font.color.rgb=RGBColor(int(text_color[-3]),int(text_color[-2]),int(text_color[-1]))
            if text_color[-3]==255 and text_color[-2]==0 and text_color[-1]==0:#赤
                for pg2 in new_zukei.text_frame.paragraphs:
                    for run2 in pg2.runs:
                        run2.font.color.rgb=RGBColor(231,75,60)
            if text_color[-3]==255 and text_color[-2]==192 and text_color[-1]==0:#橙
                for pg2 in new_zukei.text_frame.paragraphs:
                    for run2 in pg2.runs:
                        run2.font.color.rgb=RGBColor(243,157,19)
            if text_color[-3]==255 and text_color[-2]==255 and text_color[-1]==0:#黄
                for pg2 in new_zukei.text_frame.paragraphs:
                    for run2 in pg2.runs:
                        run2.font.color.rgb=RGBColor(222,176,16)
            if text_color[-3]==146 and text_color[-2]==208 and text_color[-1]==80:#薄い緑
                for pg2 in new_zukei.text_frame.paragraphs:
                    for run2 in pg2.runs:
                        run2.font.color.rgb=RGBColor(116,165,64)
            if text_color[-3]==0 and text_color[-2]==176 and text_color[-1]==80:#緑
                for pg2 in new_zukei.text_frame.paragraphs:
                    for run2 in pg2.runs:
                        run2.font.color.rgb=RGBColor(79,130,24)
            if text_color[-3]==0 and text_color[-2]==176 and text_color[-1]==240:#薄い青
                for pg2 in new_zukei.text_frame.paragraphs:
                    for run2 in pg2.runs:
                        run2.font.color.rgb=RGBColor(43,128,185)
            if text_color[-3]==0 and text_color[-2]==112 and text_color[-1]==192:#青
                for pg2 in new_zukei.text_frame.paragraphs:
                    for run2 in pg2.runs:
                        run2.font.color.rgb=RGBColor(17,62,114)
        except AttributeError:
            try:#「テーマの色」を認識
                text_color=run.font.color.theme_color
                text_color_bright=run.font.color.brightness
                #print(text_color,text_color_bright)
                for pg2 in new_zukei.text_frame.paragraphs:
                    for run2 in pg2.runs:
                        run2.font.color.theme_color=text_color
                        run2.font.color.brightness=text_color_bright
            except AttributeError:
                pass

    def change_image(self,filepass,attr,OUT_DIR):
        PPT_NAME = filepass#画像にしたいPowerPointのpath

        self.export_img(PPT_NAME, OUT_DIR)
        self.rename_img(OUT_DIR,attr)

    def export_img(self,fname, odir):
        """
        PPT_NAMEのスライドを画像に変換し、OUT_DIRに保存
        """
        pp_app = win32com.client.DispatchEx("Powerpoint.Application")
        current_folder = os.getcwd()
        try:
            shutil.rmtree(odir)
        except:
            pass
        presentation = pp_app.Presentations.open(os.path.join(current_folder, fname))

        export_path = os.path.join(current_folder, odir)
        presentation.Export(export_path, FilterName="jpg")
        presentation.close()
        pp_app.quit()

    def rename_img(self,odir,attr):
        """
        OUT_DIR内の画像の名前を「slides」に変更
        """
        file_list = glob.glob(os.path.join(odir, "*.JPG"))
        for fname in file_list:
            new_fname = fname.replace('スライド', 'slides_'+attr).lower()
            os.rename(fname, new_fname)
            if attr =='befor':
                filenamebefor.append(new_fname)
            elif attr=='after':
                filenameafter.append(new_fname)

    def get_concat_v(self,im1, im2):
        dst = Image.new('RGB', (im1.width, im1.height + im2.height))
        dst.paste(im1, (0, 0))
        dst.paste(im2, (0, im1.height))
        return dst

    def add_margin(self,pil_img, top, right, bottom, left, color):
        width, height = pil_img.size
        new_width = width + right + left
        new_height = height + top + bottom
        result = Image.new(pil_img.mode, (new_width, new_height), color)
        result.paste(pil_img, (left, top))
        return result

    def pywin(self):#画面切り替え------------------------------------
        global eff,pp_app
        pp_app = win32com.client.DispatchEx("Powerpoint.Application")
        pp_app.Visible=True
        presentation = pp_app.Presentations.open(file_1[0])
        n=presentation.Slides.Count
        for slide_number in range(1,n+1):
            if str(slide_number) in num:
                continue
            Entry_value=pp_app.ActivePresentation.Slides(slide_number).SlideShowTransition.EntryEffect
            if not Entry_value in EntryEffect:
                EntryEffect.append(Entry_value)
                if Entry_value in WrongEffect:
                    efnum.append(slide_number)
            if ef==1:
                pp_app.ActivePresentation.Slides(slide_number).SlideShowTransition.EntryEffect=0
        for i,effect in enumerate(EntryEffect,1):
            if str(i) in num:
                continue
            if effect in WrongEffect:
                eff=1
        
        presentation.SaveAs(FileName=file_22)
        presentation.close()
        pp_app.quit()

    def show_pallets(self,palette):
        fig = plt.figure(figsize=(5, 10))
        for i, color in enumerate(palette, 1):
            color_img = np.full((1, 10, 3), color, dtype=np.uint8)

            ax = fig.add_subplot(len(palette), 1, i)
            ax.imshow(color_img, aspect="auto")
            ax.set_axis_off()
            ax.text(-1, 0, i, va="center", ha="right", fontsize=10)

        plt.show()

    def iden_pallets(self):
        global background
        background=0
        for t,i in enumerate(filenameafter):
            
            s=i            
            source_file=s
            source=PIL.Image.open(source_file)
            img_rgb=source.resize((100,100))
            # Lab 色空間に変換する。
            img_lab = color.rgb2lab(img_rgb)
            palette_lab = color.rgb2lab(palette_rgb)
            # 色差を計算する。
            diff = color.deltaE_ciede2000(np.expand_dims(img_lab, axis=2), palette_lab.reshape(1, 1, -1, 3))
            
            # 一番近い色のインデックス(listの中で何番目か)を求める。
            indices = diff.argmin(axis=-1)
            #numpy.set_printoptions(threshold=numpy.inf)#省略せずに表示するのに使う
            #print(indices)
            all_color=[]
            sort_color=[]
            for i in range(len(indices)):
                nolist_indices=list(indices[i])
                for n in range(len(nolist_indices)):
                    nolist_all_color=nolist_indices[n]
            
                    all_color.append(nolist_all_color)
            #print(len(all_color))#最初のほうに100×100にしてるから10000になるはず
            for index in range(9):#0~8
                #print(all_color.count(index))
                #以下３行、色インデックスの数が少ない時、それを無視する
                if all_color.count(index) < 200:
                    while index in all_color:
                        all_color.remove(index)
            for i in range(len(all_color)):
                nolist2_all_color=all_color[i]
                if not nolist2_all_color in sort_color:
                    sort_color.append(nolist2_all_color)
            
            print('約',len(sort_color),'種類の色が使われています')
            #print(mode_color)#確認用
            s=''
            color_arr = np.array(img_rgb)
            w_size, h_size, n_color = color_arr.shape
            color_arr = color_arr.reshape(w_size * h_size, n_color)
            color_code = ['{:02x}{:02x}{:02x}'.format(*elem) for elem in color_arr]
            mode, _ = scipy.stats.mode(color_code)
            r = int(mode[0][0:2], 16)
            g = int(mode[0][2:4], 16)
            b = int(mode[0][4:6], 16)
            color_mode = (r, g, b)
            print(color_mode)#最頻値のRGB 確認用
            if color_mode!=(255,255,255):
                background=1
                sld=prs.slides[t]
                background = sld.background
                fill = background.fill
                fill.solid()
                fill.fore_color.rgb = RGBColor(255, 255, 255)#背景白にする
                backnum.append(t)
                for i,shp in enumerate(sld.shapes,1):
                    if i in num:
                        continue
                    if not shp.has_text_frame:
                        continue
                    if 'プレースホルダー' in shp.name or 'タイトル' in shp.name or 'テキスト ボックス' in shp.name:
                        for pg in shp.text_frame.paragraphs:
                            for run in pg.runs:
                                run.font.color.rgb=RGBColor(0,0,0)#文字黒にする


    def dia(self):#修正内容
        global prs
        prs = Presentation(file_22)
        
        
                #見やすいstyle_id（４行目）
        if lab==1:
            for i,slide in enumerate(prs.slides,1):
                if str(i) in num:
                    continue
                for shape in slide.shapes:
                    if not shape.has_table:
                        continue
                    tablee = shape.table
                    tbl=tablee._graphic_frame.element.graphic.graphicData.tbl
                    style_id =tbl[0][-1].text
                    #print(style_id)#表のstyle_idを取得（確認用

                    if lbl=='カラー':
                        new_style_id1='{7E9639D4-E3E2-4D34-9284-5A2195B3D0D7}'
                        new_style_id2='{69012ECD-51FC-41F1-AA8D-1B2483CD663E}'
                        new_style_id3='{72833802-FEF1-4C79-8D5D-14CF1EAF98D9}'
                        new_style_id4='{F2DE63D5-997A-4646-A377-4702673A728D}'
                        new_style_id5='{17292A2E-F333-43FB-9621-5CBBE7FDCDCB}'
                        new_style_id6='{5A111915-BE36-4E01-A7E5-04B1672EAD32}'
                        new_style_id7='{912C8C85-51F0-491E-9774-3900AFEF0FD7}'
                        #見やすいstyle_id（４行目）
                        
                        if '{2D5ABB26-0587-4C30-8999-92F81FD0307C}' in style_id or '{5940675A-B579-460E-94D1-54222C63F5DA}'\
                            in style_id or '{616DA210-FB5B-4158-B5E0-FEB733F419BA}' in style_id or '{D7AC3CCA-C797-4891-BE02-D94E43425B78}'\
                                in style_id or '{E8034E78-7F5D-4C2E-B375-FC64B27BC917}' in style_id:
                                    tbl[0][-1].text=new_style_id1
                        
                        if '{3C2FFA5D-87B4-456A-9821-1D502468CF0F}' in style_id or '{D113A9D2-9D6B-4929-AA2D-F23B5EE8CBE7}'\
                            in style_id or '{BC89EF96-8CEA-46FF-86C4-4CE0E7609802}' in style_id or '{69CF1AB2-1976-4502-BF36-3FF5EA218861}'\
                                in style_id or '{125E5076-3810-47DD-B79F-674D7AD40C01}' in style_id:
                                    tbl[0][-1].text=new_style_id2
                        
                        if '{284E427A-3D55-4303-BF80-6455036E1DE7}' in style_id or '{18603FDC-E32A-4AB5-989C-0864C3EAD2B8}'\
                            in style_id or '{5DA37D80-6434-44D0-A028-1B22A696006F}' in style_id or '{8A107856-5554-42FB-B03E-39F5DBC370BA}'\
                                in style_id or '{37CE84F3-28C3-443E-9E96-99CF82512B78}' in style_id:
                                    tbl[0][-1].text=new_style_id3
                        
                        if '{69C7853C-536D-4A76-A0AE-DD22124D55A5}' in style_id or '{306799F8-075E-4A3A-A7F6-7FBC6576F1A4}'\
                            in style_id or '{8799B23B-EC83-4686-B30A-512413B5E67A}' in style_id or '{0505E3EF-67EA-436B-97B2-0124C06EBD24}'\
                                in style_id or '{D03447BB-5D67-496B-8E87-E561075AD55C}' in style_id:
                                    tbl[0][-1].text=new_style_id4
                        
                        if '{775DCB02-9BB8-47FD-8907-85C794F793BA}' in style_id or '{E269D01E-BC32-4049-B463-5C60D7B0CCD2}'\
                            in style_id or '{ED083AE6-46FA-4A59-8FB0-9F97EB10719F}' in style_id or '{C4B1156A-380E-4F78-BDF5-A606A8083BF9}'\
                                in style_id or '{E929F9F4-4A8F-4326-A1B4-22849713DDAB}' in style_id:
                                    tbl[0][-1].text=new_style_id5
                        
                        if '{35758FB7-9AC5-4552-8A53-C91805E547FA}' in style_id or '{327F97BB-C833-4FB7-BDE5-3F7075034690}'\
                            in style_id or '{BDBED569-4797-4DF1-A0F4-6AAB3CD982D8}' in style_id or '{22838BEF-8BB2-4498-84A7-C5851F593DF1}'\
                                in style_id or '{8FD4443E-F989-4FC4-A0C8-D5A2AF1F390B}' in style_id:
                                    tbl[0][-1].text=new_style_id6
                        
                        if '{08FB837D-C827-4EFA-A057-4D05807E0F7C}' in style_id or '{638B1855-1B75-4FBE-930C-398BA8C253C6}'\
                            in style_id or '{E8B1032C-EA38-4F05-BA0D-38AFFFC7BED3}' in style_id or '{16D9F66E-5EB9-4882-86FB-DCBF35E3C3E4}'\
                                in style_id or '{AF606853-7671-496A-8E4F-DF71F8EC918B}' in style_id:
                                    tbl[0][-1].text=new_style_id7
                        #一番上と一番左のセルの文字を、セルの中心へ  

                    else:
                        new_style_id1='{8EC20E35-A176-4012-BC5E-935CFFF8708E}'
                        new_style_id2='{6E25E649-3F16-4E02-A733-19D2CDBF48F0}'
                        new_style_id3='{85BE263C-DBD7-4A20-BB59-AAB30ACAA65A}'
                        new_style_id4='{EB344D84-9AFB-497E-A393-DC336BA19D2E}'
                        new_style_id5='{EB9631B5-78F2-41C9-869B-9F39066F8104}'
                        new_style_id6='{74C1A8A3-306A-4EB7-A6B1-4F7E0EB9C5D6}'
                        new_style_id7='{2A488322-F2BA-4B5B-9748-0D474271808F}'
                        #8行目に変更する
                        
                        if '{2D5ABB26-0587-4C30-8999-92F81FD0307C}' in style_id or '{5940675A-B579-460E-94D1-54222C63F5DA}'\
                            in style_id or '{616DA210-FB5B-4158-B5E0-FEB733F419BA}' in style_id or '{D7AC3CCA-C797-4891-BE02-D94E43425B78}'\
                                in style_id or '{E8034E78-7F5D-4C2E-B375-FC64B27BC917}' in style_id:
                                    tbl[0][-1].text=new_style_id1
                        
                        if '{3C2FFA5D-87B4-456A-9821-1D502468CF0F}' in style_id or '{D113A9D2-9D6B-4929-AA2D-F23B5EE8CBE7}'\
                            in style_id or '{BC89EF96-8CEA-46FF-86C4-4CE0E7609802}' in style_id or '{69CF1AB2-1976-4502-BF36-3FF5EA218861}'\
                                in style_id or '{125E5076-3810-47DD-B79F-674D7AD40C01}' in style_id:
                                    tbl[0][-1].text=new_style_id2
                        
                        if '{284E427A-3D55-4303-BF80-6455036E1DE7}' in style_id or '{18603FDC-E32A-4AB5-989C-0864C3EAD2B8}'\
                            in style_id or '{5DA37D80-6434-44D0-A028-1B22A696006F}' in style_id or '{8A107856-5554-42FB-B03E-39F5DBC370BA}'\
                                in style_id or '{37CE84F3-28C3-443E-9E96-99CF82512B78}' in style_id:
                                    tbl[0][-1].text=new_style_id3
                        
                        if '{69C7853C-536D-4A76-A0AE-DD22124D55A5}' in style_id or '{306799F8-075E-4A3A-A7F6-7FBC6576F1A4}'\
                            in style_id or '{8799B23B-EC83-4686-B30A-512413B5E67A}' in style_id or '{0505E3EF-67EA-436B-97B2-0124C06EBD24}'\
                                in style_id or '{D03447BB-5D67-496B-8E87-E561075AD55C}' in style_id:
                                    tbl[0][-1].text=new_style_id4
                        
                        if '{775DCB02-9BB8-47FD-8907-85C794F793BA}' in style_id or '{E269D01E-BC32-4049-B463-5C60D7B0CCD2}'\
                            in style_id or '{ED083AE6-46FA-4A59-8FB0-9F97EB10719F}' in style_id or '{C4B1156A-380E-4F78-BDF5-A606A8083BF9}'\
                                in style_id or '{E929F9F4-4A8F-4326-A1B4-22849713DDAB}' in style_id:
                                    tbl[0][-1].text=new_style_id5
                        
                        if '{35758FB7-9AC5-4552-8A53-C91805E547FA}' in style_id or '{327F97BB-C833-4FB7-BDE5-3F7075034690}'\
                            in style_id or '{BDBED569-4797-4DF1-A0F4-6AAB3CD982D8}' in style_id or '{22838BEF-8BB2-4498-84A7-C5851F593DF1}'\
                                in style_id or '{8FD4443E-F989-4FC4-A0C8-D5A2AF1F390B}' in style_id:
                                    tbl[0][-1].text=new_style_id6
                        
                        if '{08FB837D-C827-4EFA-A057-4D05807E0F7C}' in style_id or '{638B1855-1B75-4FBE-930C-398BA8C253C6}'\
                            in style_id or '{E8B1032C-EA38-4F05-BA0D-38AFFFC7BED3}' in style_id or '{16D9F66E-5EB9-4882-86FB-DCBF35E3C3E4}'\
                                in style_id or '{AF606853-7671-496A-8E4F-DF71F8EC918B}' in style_id:
                                    tbl[0][-1].text=new_style_id7
                        #１，２、５、９、１０行目を8行目（見やすいやつ）に変更
                    
                    for c in range(len(tablee.columns)):
                        cell1 = tablee.cell(0, c)
                        cell1.vertical_anchor = MSO_ANCHOR.MIDDLE
                        pg = cell1.text_frame.paragraphs[0]
                        pg.alignment = PP_ALIGN.CENTER
                    for r in range(len(tablee.rows)):
                        cell2 = tablee.cell(r, 0)
                        cell2.vertical_anchor=MSO_ANCHOR.MIDDLE
                        pg = cell2.text_frame.paragraphs[0]
                        pg.alignment = PP_ALIGN.CENTER
                    #一番上と一番左のセルの文字を、セルの中心へ
                        
                    for x in range(len(tablee.columns)):
                        for y in range(len(tablee.rows)):
                            cell3 = tablee.cell(y, x)
                            cell3.fill.solid()
                            cell3.fill.fore_color.rgb=RGBColor(255,255,255)#文字を黒に統一
                            pg = cell3.text_frame.paragraphs[0]
                            pg.font.color.rgb=RGBColor(0,0,0)#セルをすべて白で塗りつぶし
                            cell_text=pg.text
                            
                            if cell_text.isdigit():#もし数字なら
                                cell3.vertical_anchor=MSO_ANCHOR.MIDDLE
                                pg.alignment = PP_ALIGN.RIGHT#数字は右揃え＆上下中央揃え
                            
                            for run in pg.runs:
                                run.font.italic=None
                                run.font.underline=None#イタリックと下線をなくす

        self.displacement() #displacementを実行する
        #ご静聴削除
        for i,slide in enumerate(prs.slides,1):
            if str(i) in num:
                continue
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                if"清聴" in shape.text or '静聴'in shape.text: 
                    xml_slides = prs.slides._sldIdLst
                    slides = list(xml_slides)
                    xml_slides.remove(slides[-1])

                #左揃え
                if lef==1:
                    if 'プレースホルダー' in shape.name or 'タイトル' in shape.name or 'テキスト ボックス' in shape.name:
                        for pg in shape.text_frame.paragraphs:
                            if pg.alignment!=1:
                                pg.alignment = 1             
        """if self.sw.active==True:
            sm.current = 'edu'"""
        self.change_image(file_22,'after',OUT_DIR_after)
        self.Shape_change()
        self.iden_pallets()
        prs.save(file_22)
        self.display(sent)


        self.initialization()

    def initialization(self):
        global file_1,file_2,file2,prs,eff,sld1,sent,EntryEffect,\
                shp_name,plc_l,plc_w,plc_t,plc_h,sd,shap,bold,ita,und,\
                    para,fonts,texts,runs,color_1,color_2,color_3,color_t,\
                        color_b,center_x,center_y,right,bottom,rotation,df,talbe_tex,shpname,\
                            filenamebefor,filenameafter,background,ef,it,un,fnt,lef,fontname,font,lo,\
                                line,cir,ex,wa,scr,arr,lbl,lab,efe,num,path,italnum,boldnum,lenghnum,efnum,backnum,da
                            
        file_1='aaaaa'
        file_2='aaaaa'
        file2=''
        font='メイリオ'
        prs=None
        eff=0
        ef=1
        it=1
        un=1
        fnt=1
        lef=1
        line=1
        cir=1
        ex=1
        wa=1
        scr=1
        arr=1
        lab=1
        efe=1
        lo=0
        path=0
        da=0
        sld1=None
        background=0
        sent='保存しました\n'
        EntryEffect=[]
        shp_name=[]
        plc_l=[]
        plc_w=[]
        plc_t=[]
        plc_h=[]
        sd=[]
        shap=[]
        bold=[]
        ita=[]
        und=[]
        fontname=[]
        filenamebefor=[]
        filenameafter=[]
        para=[]
        fonts=[]
        texts=[]
        runs=[]
        color_1=[]
        color_2=[]
        color_3=[]
        color_t=[]
        color_b=[]
        center_x=[]
        center_y=[]
        right=[]
        bottom=[]
        rotation=[]
        df=pd.DataFrame({})
        talbe_tex=[]
        shpname=[]
        num=[]
        lbl='カラー'
        PopupMenue_setting().sw1.active=True
        PopupMenue_setting().sw2.active=True
        PopupMenue_setting().sw3.active=True
        PopupMenue_setting().sw4.active=True
        PopupMenue_setting().sw5.active=True
        PopupMenue_setting().sw6.active=True
        PopupMenue_setting().sw10.active=False
        PopupMenue_setting().sw8.active=True
        PopupMenue_setting().sw9.active=True
        italnum=[]
        boldnum=[]
        lenghnum=[]
        efnum=[]
        backnum=[]

    def display(self,sent):
        global da
        if PopupMenue_setting().sw4.active==True:
            if italnum!=None:
                sent=sent+'\nイタリックはやめましょう'
        if PopupMenue_setting().sw3.active==True:
            if boldnum!=None:
                sent=sent+'\n下線はやめましょう'
        if PopupMenue_setting().sw2.active==True:
            if efnum!=None and ef==1:
                sent=sent+'\n派手な画面切り替えは控えましょう'
        if backnum!=1:
            sent=sent+'\n白色ベースにしましょう'
        if lenghnum!=1:
            sent=sent+'\n文章は簡潔にしましょう'
        if sent!='保存しました\n':
            self.lbl.text=sent
        else:
            self.lbl.text="アドバイスはありませんでした！"




    def popup_open(self):
        content = PopupMenu(popup_close=self.popup_close)
        self.popup = Popup(title='選択', content=content, size_hint=(0.75, 0.7), auto_dismiss=False)
        self.popup.open()
    #ポップアップを開く

    def popup_open2(self):
        content = PopupMenu2(popup_close2=self.popup_close2)
        self.popup = Popup(title='保存', content=content, size_hint=(0.75, 0.7), auto_dismiss=False)
        if file_1!='aaaaa' and file_1!=[]:
            self.popup.open()
            PopupMenu2().refilename()
        else:
            self.lbl.text='ファイルを保存してください'
    #ポップアップを開く
    def setting_open(self):
        content = PopupMenue_setting(popup_close_setting=self.popup_close_setting)
        self.popup = Popup(title='設定', content=content, size_hint=(0.75, 0.7), auto_dismiss=False)
        self.popup.open()

    def popup_close_setting(self):
        self.popup.dismiss()

    def popup_close(self):
        global file_1,path
        self.popup.dismiss()
        if file_1!='aaaaa' and file_1!=[]:
            self.lbl.text='名前を付けて保存してください' 
            self.change_image(file_1[0],"befor",OUT_DIR_befor)
            
            prs = Presentation(file_1[0])
            n=len(prs.slides)

            im1 = Image.open('python-pptx_images_befor\\slides_befor1.jpg')
            Height=im1.height
            for i in range(n-1):
                
                page_number=str(i+2)#二枚目から
                im=Image.open('python-pptx_images_befor\\slides_befor'+ page_number +'.jpg')
                im1=self.get_concat_v(im1, im)

            im_new = self.add_margin(im1, 0, 0, 0, 600, (77, 77, 77))
            font_path = 'YuGothB.ttc'
            font_size=500
            font = ImageFont.truetype(font_path, font_size)
            for i in range(n):
                page_number=str(i+1)
                draw = ImageDraw.Draw(im_new)
                location= i * Height
                draw.text((10,location),page_number, '#FFFFFF' , font=font)

            im_new.save('python-pptx_images_befor\\binded.jpg')
            path=1
            if lo == 1:
                sm.transition.direction='right'#画面切り替えの時右に動かす
                sm.current = 'lock'




    #ポップアップを閉じる
    def popup_close2(self):
        global file_22
        self.popup.dismiss()
        if file2=='':
            return
        file_22='{}\{}.pptx'.format(file_2,file2)
        self.pywin()
        self.dia()
        
    #ポップアップを閉じる
    def open_instruction(self): 
        sm.current = 'Instructions' #select画面にする
        sm.transition.direction='left' #画面変更の時左に移動する

    def pressedu(self):
        sm.transition.direction='left'
        sm.current = 'edu'

    def displacement(self):
        global ital,unde,lenth
        ital=0
        unde=0
        lenth=0
        cnt=-1
        for sld in prs.slides:
            cnt+=1
            count=-1
            for shp in sld.shapes:
                cont=-1
                count+=1
                ll=shp.left.pt #左端の位置を取得する
                tt=shp.top.pt#頂上の位置を取得する
                ww=shp.width.pt#幅を取得する
                hh=shp.height.pt#高さを取得する
                shame=shp.name
                if not shp.has_text_frame:
                    continue
                for pg1 in shp.text_frame.paragraphs:
                    cont+=1
                    run=-1
                    for run1 in pg1.runs:
                        run+=1
                        k=run1.font.size
                        t=run1.text
                        bl=run1.font.bold
                        it=run1.font.italic
                        un=run1.font.underline
                        fname=run1.font.name
                        try:
                            text_color=run1.font.color.rgb
                            color_1.append(int(text_color[0]))
                            color_2.append(int(text_color[1]))
                            color_3.append(int(text_color[2]))
                            text_color_t=-1
                            text_color_bright=-1
                        except AttributeError:
                            try:
                                text_color_t=run1.font.color.theme_color
                                text_color_bright=run1.font.color.brightness
                                color_1.append(-1)
                                color_2.append(-1)
                                color_3.append(-1)
                            except AttributeError:
                                color_1.append(0)
                                color_2.append(0)
                                color_3.append(0)
                                text_color_t=-1
                                text_color_bright=-1
        
                        if k!=None:
                            fig=int(k/12700)
                        else:
                            fig=None
                        rotation.append(shp.rotation)
                        sd.append(cnt)
                        shap.append(count)
                        para.append(cont)
                        runs.append(run)
                        shpname.append(shame)

                        plc_h.append(hh)
                        plc_w.append(ww)
                        plc_t.append(tt)
                        plc_l.append(ll)
                        fonts.append(fig)
                        texts.append(t)
                        bold.append(bl)
                        ita.append(it)
                        und.append(un)
                        fontname.append(fname)
                        color_t.append(text_color_t)
                        color_b.append(text_color_bright)
                        #それぞれの値をリストに入れる

        df['slide']=sd
        df['shape']=shap
        df['paragraph']=para
        df['run']=runs
        df['size']=fonts
        df['text']=texts
        df['shp_name']=shpname
        df['Left']=plc_l
        df['Top']=plc_t
        df['width']=plc_w
        df['Height']=plc_h
        df['Rotation']=rotation
        df['Color_G']=color_2
        df['Color_B']=color_3
        df['T_Color']=color_t
        df['B_Color']=color_b
        #それぞれのリストをデータフレームに入れる
        df.fillna({'size':-5},inplace=True)#sizeのNONEのところを‐5に置き換える
        for i in df.index:
            center_x.append(df.at[i,'Left']+df.at[i,'width']/2)#ｘ軸の中心位置を計算するリストに入れる
            center_y.append(df.at[i,'Top']+df.at[i,'Height']/2)#ｙ軸の中心位置を計算してリストに入れる
            right.append(df.at[i,'Left']+df.at[i,'width'])#右端の位置を計算してリストに入れる
            bottom.append(df.at[i,'Top']+df.at[i,'Height'])#底の位置を計算してリストに入れる
        df['center_x']=center_x
        df['center_y']=center_y
        df['Right']=right
        df['Bottom']=bottom
        df['bold']=bold
        df['italic']=ita
        df['underline']=und
        df['font_name']=fontname
        df['Color_R']=color_1
        

        #テキストボックス内の文字の引継ぎ、修正--------------------------------
        
        for i,sld1 in enumerate(prs.slides,1):
            if str(i) in num:
                continue

            for shp1 in sld1.shapes:
                if 'プレースホルダー' in shp1.name or 'タイトル' in shp1.name or 'テキスト ボックス' in shp1.name:
                    if not shp1.has_text_frame:
                        continue
                    for pg1 in shp1.text_frame.paragraphs:
                        pg1.text=''
                if line==1:
                    try:#標準の色
                        shape_color=shp1.fill.fore_color.rgb
                        colo=RGBColor(int(shape_color[-3]),int(shape_color[-2]),int(shape_color[-1]))
                
                        if not shape_color[-3]==255 or not shape_color[-2]==255 or not shape_color[-1]==255:
                            shp1.line.color.rgb=RGBColor(255,255,255)
                    
                    except AttributeError:
                        try:#テーマの色
                            shape_color=shp1.fill.fore_color.theme_color
                            shape_color_bright=shp1.fill.fore_color.brightness
                            if shape_color==14 and shape_color_bright > -0.15:
                                continue
                            shp1.line.color.rgb=RGBColor(255,255,255)
                        except AttributeError:
                            pass
                    
                    
                    except TypeError as e:#塗りつぶしが初期色or塗りつぶしなしorグラデーションとかいろいろ

                        print('＝＝＝＝＝エラー内容＝＝＝＝＝')
                        print('type:' + str(type(e)))
                        print(str(e))
                        
                        if '_NoFill has no foreground color' in str(e):#「塗りつぶしなし」のときのメッセージ
                            print('これは塗りつぶしなしです')
                        if '_NoneFill has no foreground color' in str(e):#「塗りつぶし」が初期色のときのメッセージ
                            print('塗りつぶしが初期色です')
                        if '_GradFill has no foreground color' in str(e):#「塗りつぶし」がグラデーションのときのメッセージ
                            print('塗りつぶしがグラデーションです')
                        if '_BlipFill has no foreground color' in str(e):#「塗りつぶし」がテクスチャのときのメッセージ
                            print('塗りつぶしがテクスチャです')

        for index in df.index:
            print(index)
            if str(df.at[index,'slide']+1) in num:
                continue
            if df.at[index,'italic']==True:
                italnum.append(df.at[index,'slide'])
                
            if df.at[index,'underline']==True:
                boldnum.append(df.at[index,'slide'])
            sld2=prs.slides[df.at[index,'slide']]
            shp2=sld2.shapes[df.at[index,'shape']]
            pg2=shp2.text_frame.paragraphs[df.at[index,'paragraph']]
            new_text=pg2.add_run()
            if it==1:
                new_text.font.italic=None
            else:
                new_text.font.italic=df.at[index,'italic']
            if un==1:
                new_text.font.underline=None
            else:
                new_text.font.underline=df.at[index,'underline']
            new_text.text=df.at[index,'text']
            if fnt == 1:
                new_text.font.name=font
            else:
                new_text.font.name=df.at[index,'font_name']
            new_text.font.bold=df.at[index,'bold']
            
            if df.at[index,'T_Color']==-1:
                new_text.font.color.rgb=RGBColor(int(df.at[index,'Color_R']),int(df.at[index,'Color_G']),int(df.at[index,'Color_B']))
            else:
                new_text.font.color.theme_color=df.at[index,'T_Color']
                new_text.font.color.brightness=df.at[index,'B_Color']
            if df.at[index,'size']==-5:
                continue
            new_text.font.size=Pt(df.at[index,'size'])
        for slide in range(0,df.iloc[-1]['slide']) :
            length=0
            for index in df.index:
                if  df.at[index,'slide']==slide:
                    text=df.at[index,'text']
                    length=length+len(text)
            if length>105:
                lenghnum.append(df.at[index,'slide'])

        
        
        for i,sld in enumerate(prs.slides,1):
            if str(i) in num:
                continue
            for shp_k in sld.shapes:
                if 'タイトル' in shp_k.name:
                    continue
                
                l_k=shp_k.left
                ll_k=int(l_k)/12700
                lll_k=Pt(ll_k)
                t_k=shp_k.top
                tt_k=int(t_k)/12700
                ttt_k=Pt(tt_k)
        
                for shp_i in sld.shapes:
                    if 'タイトル' in shp_i.name:
                        continue

                    l_i=shp_i.left
                    ll_i=int(l_i)/12700
                    lll_i=Pt(ll_i)
                    
                    t_i=shp_i.top
                    tt_i=int(t_i)/12700
                    ttt_i=Pt(tt_i)
                    
                    if -500000 < lll_i - lll_k <0:
                        shp_k.left=lll_i
                    if 0 < lll_i - lll_k <500000:
                        shp_i.left=lll_k
                    if -500000 < ttt_i - ttt_k <0:
                        shp_k.top=ttt_i
                    if 0 < ttt_i - ttt_k <500000:
                        shp_i.top=ttt_k
        
                
        
            
                #テキストボックスと図形が何かの色で塗りつぶされてるとき「枠線」を白にする-------
        for i,sld in enumerate(prs.slides,1):
            if str(i) in num:
                continue
            for shp in sld.shapes:
                try:#標準の色
                    shape_color=shp.fill.fore_color.rgb
                    colo=RGBColor(int(shape_color[-3]),int(shape_color[-2]),int(shape_color[-1]))
        
                    if not shape_color[-3]==255 or not shape_color[-2]==255 or not shape_color[-1]==255:
                        shp.line.color.rgb=RGBColor(255,255,255)

                except AttributeError:
                    try:#テーマの色
                        shape_color=shp.fill.fore_color.theme_color
                        shape_color_bright=shp.fill.fore_color.brightness
                        if shape_color==14 and shape_color_bright > -0.15:
                            continue
                        shp.line.color.rgb=RGBColor(255,255,255)
                    except AttributeError:
                        pass

                except TypeError as e:#塗りつぶしが初期色or塗りつぶしなしorグラデーションとかいろいろ
                    '''
                    print('＝＝＝＝＝エラー内容＝＝＝＝＝')
                    print('type:' + str(type(e)))
                    print(str(e))
                    '''
                    if '_NoneFill has no foreground color' in str(e):#「塗りつぶし」が初期色のとき
                        if 'プレースホルダー' in shp.name or 'タイトル' in shp.name or 'テキスト ボックス' in shp.name:
                            pass
                        else:
                            shp.line.color.rgb=RGBColor(255,255,255)
                    elif '_NoFill has no foreground color' in str(e):#「塗りつぶしなし」のとき
                        pass
                    else:
                        shp.line.color.rgb=RGBColor(255,255,255)
            #アドバイス表示(display関数)
    def Shape_change(self):
        gray=RGBColor(160,160,160)
        for i,sld in enumerate(prs.slides,1):
            if str(i) in num:
                continue
            shapes=sld.shapes
            for shp in sld.shapes:
                if efe==1:
                    if shp.has_table:
                        continue
                    """shp.shadow.inherit=True"""
                l=shp.left
                ll=int(l)/12700
                lll=Pt(ll)
                t=shp.top
                tt=int(t)/12700
                ttt=Pt(tt)
                w=shp.width
                ww=int(w)/12700
                www=Pt(ww)
                h=shp.height
                hh=int(h)/12700
                hhh=Pt(hh)
                if hasattr(shp, 'text'): 
                    s=shp.text
                f=''
                degree=shp.rotation
                if scr==1:
                    if  'スクロール'in shp.name or '巻き' in shp.name:
                        ww_adjust=int(ww)*92/100#横の長さを少し短くする調整
                        www=Pt(ww_adjust)
                        new_zukei = shapes.add_shape(MSO_SHAPE.RECTANGLE,lll,ttt,www,hhh)
                        new_zukei.text=s#文字を引き継ぐ
                        for pg in shp.text_frame.paragraphs:
                            for run in pg.runs:
                                f=run.font.size
                                if not run.font.size==None:
                                    ff=int(f)/12700
                                    fff=Pt(ff)
                                else:
                                    fff=Pt(18)
                                self.text_color_pass(run,new_zukei)#文字色を引き継ぐ関数
                        self.shape_color_pass(shp,new_zukei)
                        for pg2 in new_zukei.text_frame.paragraphs:
                            pg2.alignment=PP_ALIGN.CENTER#文字を中央揃え
                            for run2 in pg2.runs:
                                run2.font.size=fff#文字サイズを引き継ぐ
                        shp.width=0
                        shp.height=0
                        if  shp.has_text_frame:
                            shp.text_frame.clear()
        # --------------------爆発-------------------- 
                if ex==1:
                    if  '爆発' in shp.name:
                        ww_adjust=int(ww)*68/100#ここから調整
                        hh_adjust=int(hh)*47/100
                        shaft_h=int(hh)/2 + int(tt)
                        tt=int(shaft_h)-int(hh_adjust)/2
                        shaft_w=int(ww)/2 + int(ll)
                        ll=int(shaft_w)-int(ww_adjust)/2                    
                        www=Pt(ww_adjust)
                        hhh=Pt(hh_adjust)
                        new_zukei = shapes.add_shape(MSO_SHAPE.RECTANGLE,lll,ttt,www,hhh)
                        new_zukei.text=s#文字を引き継ぐ
                        for pg in shp.text_frame.paragraphs:
                            for run in pg.runs:
                                f=run.font.size
                                if not run.font.size==None:
                                    ff=int(f)/12700
                                    fff=Pt(ff)
                                else:
                                    fff=Pt(18)
                                self.text_color_pass(run,new_zukei)#文字色を引き継ぐ関数
                        self.shape_color_pass(shp,new_zukei)
                        for pg2 in new_zukei.text_frame.paragraphs:
                            pg2.alignment=PP_ALIGN.CENTER#文字を中央揃え
                            for run2 in pg2.runs:
                                run2.font.size=fff#文字サイズを引き継ぐ
                        shp.width=0
                        shp.height=0
                        if  shp.has_text_frame:
                            shp.text_frame.clear()
        # --------------------波形、リボン--------------------
                if wa==1:
                    if  '波'  in shp.name or  'リボン' in shp.name:
                        new_zukei = shapes.add_shape(MSO_SHAPE.RECTANGLE,lll,ttt,www,hhh)
                        new_zukei.text=s#文字を引き継ぐ
                        for pg in shp.text_frame.paragraphs:
                            for run in pg.runs:
                                f=run.font.size
                                if not run.font.size==None:
                                    ff=int(f)/12700
                                    fff=Pt(ff)
                                else:
                                    fff=Pt(18)
                                self.text_color_pass(run,new_zukei)#文字色を引き継ぐ関数
                        self.shape_color_pass(shp,new_zukei)
                        for pg2 in new_zukei.text_frame.paragraphs:
                            pg2.alignment=PP_ALIGN.CENTER#文字を中央揃え
                            for run2 in pg2.runs:
                                run2.font.size=fff#文字サイズを引き継ぐ
                        shp.width=0
                        shp.height=0
                        if  shp.has_text_frame:
                            shp.text_frame.clear()
        #　--------------------円--------------------
                if cir==1:
                    if  '楕円'  in shp.name:
                        if 1<=int(ww)/int(hh)<1.5:
                            hh=ww
                            shp.height=Pt(hh)
                        elif 1<int(hh)/int(ww)<1.5:
                            ww=hh
                            shp.width=Pt(ww)
                        else:#楕円の時は四角形にする
                            s=shp.text
                            lll=Pt(ll)
                            ttt=Pt(tt)
                            www=Pt(ww)
                            hhh=Pt(hh)
                            new_zukei = shapes.add_shape(MSO_SHAPE.RECTANGLE,lll,ttt,www,hhh)
                            new_zukei.text=s#文字を引き継ぐ
                            for pg in shp.text_frame.paragraphs:
                                for run in pg.runs:
                                    f=run.font.size
                                    if not run.font.size==None:
                                        ff=int(f)/12700
                                        fff=Pt(ff)
                                    else:
                                        fff=Pt(18)
                                    self.text_color_pass(run,new_zukei)#文字色を引き継ぐ関数
                            self.shape_color_pass(shp,new_zukei)
                            for pg2 in new_zukei.text_frame.paragraphs:
                                pg2.alignment=PP_ALIGN.CENTER#文字を中央揃え
                                for run2 in pg2.runs:
                                    run2.font.size=fff#文字サイズを引き継ぐ
                            shp.width=0
                            shp.height=0
                            if  shp.has_text_frame:
                                shp.text_frame.clear()
        #矢印--------------------------------
                if arr==1:
                    for arro in arrow:
                        if arro in shp.name:
                            if ww < 40 and hh > 200:
                                shp.fill.solid()#矢印の色を灰色にする
                                shp.fill.fore_color.rgb=gray
                                shp.line.color.rgb = gray
                                continue
                            if ww > 330 and hh < 40:
                                shp.fill.solid()#矢印の色を灰色にする
                                shp.fill.fore_color.rgb=gray
                                shp.line.color.rgb = gray
                                continue
                            shaft_h=int(hh)/2 + int(tt)#ここから調整
                            shaft_w=int(ww)/2 + int(ll)
                            if int(hh)/int(ww)>1.5:
                                hh=ww
                                tt=int(shaft_h)-int(hh)/2
                            if int(ww)/int(hh)>1.5:
                                ww=hh
                                ll=int(shaft_w)-int(ww)/2#ここまで調整
                            lll=Pt(ll)
                            ttt=Pt(tt)
                            www=Pt(ww)
                            hhh=Pt(hh)
                            
                            for pg in shp.text_frame.paragraphs:
                                for run in pg.runs:
                                    f=run.font.size
                                    if not run.font.size==None:
                                        ff=int(f)/12700
                                        fff=Pt(ff)
                                    else:
                                        fff=Pt(18)
                            if '矢印: 下' in shp.name or '下矢印' in shp.name:
                                new_zukei = shapes.add_shape(MSO_SHAPE.FLOWCHART_MERGE,lll,ttt,www,hhh)
                                new_zukei.rotation=degree
                            elif '矢印: 上' in shp.name or '上矢印' in shp.name:
                                new_zukei = shapes.add_shape(MSO_SHAPE.FLOWCHART_EXTRACT,lll,ttt,www,hhh)
                                new_zukei.rotation=degree
                            elif  '矢印: 右' in shp.name or '右矢印' in shp.name:
                                new_zukei = shapes.add_shape(MSO_SHAPE.FLOWCHART_EXTRACT,lll,ttt,www,hhh)
                                new_zukei.rotation=int(degree)+90
                            elif  '矢印: 左' in shp.name or '左矢印' in shp.name:
                                new_zukei = shapes.add_shape(MSO_SHAPE.FLOWCHART_MERGE,lll,ttt,www,hhh)
                                new_zukei.rotation=int(degree)+90
                            new_zukei.text=s
                            new_zukei.fill.solid()#矢印の色を灰色にする
                            new_zukei.fill.fore_color.rgb=gray
                            new_zukei.line.color.rgb = gray
                            for pg2 in new_zukei.text_frame.paragraphs:
                                pg2.alignment=PP_ALIGN.CENTER
                                for run2 in pg2.runs:
                                    run2.font.size=fff
                            shp.width=0
                            shp.height=0
                            if  shp.has_text_frame:
                                shp.text_frame.clear()








class rrr (Screen):
    def __init__(self, **kwargs):
        super(rrr, self).__init__(**kwargs)
        pass
    
    def back(self):
        sm.transition.direction='right'#画面切り替えの時右に動かす
        sm.current = 'menu'#MENU画面を表示する

class PopupMenu(Screen):
    popup_close = ObjectProperty(None)
    path=StringProperty('')
    def filea(self):
        global file_1,file_name
        file_1=self.vv.selection#kivyファイルのfile_1をselectionにする
        
        
        if file_1==[] :
            self.bt.text='ファイルを選択してください'
        else:
            self.popup_close()
            file_name=file_1[0].split('\\')[-1].split('.')[0] +'_修正後'

class PopupMenu2(Screen):
    popup_close2 = ObjectProperty(None)
    path =StringProperty('')

    def refilename(self):
        self.input.text=file_name
        
    def fileb(self):
        global file_1,file_2,prs,sld1,file2
        file_2=self.vv.path  
        file2=self.input.text
        if file2=='':
            self.input.text=file_name
        else:
            self.popup_close2()

class Lock(Screen):
    def __init__(self, **kwargs):
        super(Lock, self).__init__(**kwargs)
        Clock.schedule_interval(self.update, 60/60.0)
        pass

    source = StringProperty('')
    def update(self,dt):
        self.im.reload
        if path==1:
            self.source='python-pptx_images_befor//binded.jpg'
    
    def back(self):
        sm.transition.direction='left'#画面切り替えの時右に動かす
        sm.current = 'menu'#MENU画面を表示する

    def Lock_number(self):
        global num
        number_text=self.txin.text
        num=re.findall(r"\d+", number_text)
        self.lbnum.text='{}は、修正しない'.format(num)




class PopupMenue_setting(Screen):
    popup_close_setting = ObjectProperty(None)
    def setting(self):
        global fnt,ef,it,un,lef,font,line,cir,ex,wa,scr,arr,lbl,lab,efe,lo
        if self.sw1.active:
            fnt=1
        else :
            fnt=0

        if self.sw2.active:
            ef=1
        else :
            ef=0

        if self.sw3.active:
            un=1
        else :
            un=0

        if self.sw4.active:
            it=1
        else :
            it=0

        if self.sw5.active:
            lef=1
        else:
            lef=0

        if self.sw6.active:
            line=1
        else :
            line=0
        
        if self.cir.state=='down':
            cir=1
        else:
            cir=0
        
        if self.ex.state=='down':
            ex=1
        else:
            ex=0
        
        if self.wave.state=='down':
            wa=1
        else:
            wa=0
        
        if self.scr.state=='down':
            scr=1
        else:
            scr=0
        
        if self.arr.state=='down':
            arr=1
        else:
            arr=0
        
        if self.lbl.text=='カラー':
            lbl='カラー'
        else:
            lbl='白黒'
        
        if self.sw8.active:
            lab=1
        else:
            lab=0
        
        if self.sw9.active:
            efe=1
        else:
            efe=0
        
        if self.sw10.active:
            lo=1
        else:
            lo=0


        if self.lb.text=='メイリオ':
            font='メイリオ'
        else:
            font='游ゴシック'
        self.popup_close_setting()


class StudyApp(App):
    def build(self):
        sm.add_widget(diagnosis(name='menu'))#menuという名前の画面を作る
        sm.add_widget(rrr(name='Instructions'))
        sm.add_widget(Lock(name='lock'))#selectという名前の画面を作る
        self.icon = r'pic\icon.png'
        return sm 



if __name__ == '__main__':#このプログラムが直接呼び出されたとき
    file_1='aaaaa'
    file_2='aaaaa'
    sm = ScreenManager()
    file2=''
    font='メイリオ'
    lbl='カラー'
    prs=None
    eff=0
    ef=1
    it=1
    un=1
    fnt=1
    lef=1
    line=1
    cir=1
    ex=1
    wa=1
    scr=1
    arr=1
    lab=1
    efe=1
    path=0
    lo=0
    da=0

    sld1=None
    sent='保存しました\n'
    
    EntryEffect=[]
    WrongEffect=[513,3850,3863,3864,3865,
                 3866,3872,3873,3874,3875,
                 3876,3877,3878,3879,3898,
                 3899,3900,3910,3911,3912,
                 3913,3942,3943,3950,3951,
                 3952,3953]
    arrow=['矢印: 下','下矢印','矢印: 上','上矢印','矢印: 右','右矢印','矢印: 左','左矢印']
    shp_name=[]
    plc_l=[]
    plc_w=[]
    plc_t=[]
    plc_h=[]
    sd=[]
    shap=[]
    bold=[]
    ita=[]
    und=[]
    para=[]
    fonts=[]
    texts=[]
    fontname=[]
    runs=[]
    filenamebefor=[]
    filenameafter=[]
    color_1=[]
    color_2=[]
    color_3=[]
    color_t=[]
    color_b=[]
    center_x=[]
    center_y=[]
    right=[]
    bottom=[]
    rotation=[]
    df=pd.DataFrame({})
    talbe_tex=[]
    shpname=[]
    num=[]
    italnum=[]
    boldnum=[]
    lenghnum=[]
    efnum=[]
    backnum=[]

    palette_rgb = np.array(
    [
        (255,255,255),
        (255, 0, 0),
        (255, 192, 0),
        (255, 255, 0),
        (0, 176, 80),
        (0, 112, 192),
        (112, 48, 160),
        (20, 20, 20),
        (17, 60, 73),
        
    ],
    dtype=np.uint8,
)
    
    OUT_DIR_after = 'python-pptx_images_after'
    OUT_DIR_befor = 'python-pptx_images_befor'#画像を保存するフォルダ名
    pathTest = r"python-pptx_images"

    

    StudyApp().run()#StudyAppを動かす